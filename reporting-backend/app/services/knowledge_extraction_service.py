import base64
from datetime import datetime
from html.parser import HTMLParser
import io
import ipaddress
import logging
import re
import socket
from typing import Any, Dict, List
from urllib.parse import urljoin, urlparse

import pandas as pd
import requests
from requests.adapters import HTTPAdapter
import urllib3
from fastapi import HTTPException

from app.core.gemini_config import chat_fallback_chain
from app.core.supabase import supabase
from app.services.knowledge_embedding_service import (
    GEMINI_API_BASE,
    _HTTPX_CLIENT,
    _embed_texts,
    _get_gemini_keys,
    _gemini_api_key,
    _vector_literal,
)

logger = logging.getLogger(__name__)

MAX_STORAGE_UPLOAD_SIZE_BYTES = 50 * 1024 * 1024
MAX_WEB_PAGE_BYTES = 1 * 1024 * 1024
MIN_EXTRACTED_TEXT_CHARS = 20
MAX_CONTEXTUAL_CHUNKS_PER_DOCUMENT = 40
IGNORED_HTML_TAGS = {"script", "style", "noscript", "svg", "nav", "header", "footer", "aside"}

MOJIBAKE_REPLACEMENTS = {
    "\xa0": " ",
    "\u200b": "",
    "\u2022": "• ",
    "\uf0b7": "• ",
    "\u2013": "-",
    "\u2014": "-",
    "\u2018": "'",
    "\u2019": "'",
    "\u201c": '"',
    "\u201d": '"',
    "\ufeff": "",
    "\u00c2\u00b0": "\u00b0",
    "\u00c2": "",
    "\u00e2\u20ac\u00a2": "- ",
    "\u00e2\u20ac\u201c": "-",
    "\u00e2\u20ac\u201d": "-",
    "\u00e2\u20ac\u02dc": "'",
    "\u00e2\u20ac\u2122": "'",
    "\u00e2\u20ac\u0153": '"',
    "\u00e2\u20ac\ufffd": '"',
}


class HostPinnedHTTPAdapter(HTTPAdapter):
    def __init__(self, hostname: str, pinned_ip: str, **kwargs):
        self.hostname = hostname
        self.pinned_ip = pinned_ip
        super().__init__(**kwargs)

    def init_poolmanager(self, connections, maxsize, block=False, **pool_kwargs):
        class PinnedPoolManager(urllib3.PoolManager):
            def __init__(pm_self, hostname, pinned_ip, *args, **kwargs):
                pm_self.hostname = hostname
                pm_self.pinned_ip = pinned_ip
                super().__init__(*args, **kwargs)

            def _new_pool(pm_self, scheme, host, port, request_context=None):
                if host == pm_self.hostname:
                    if request_context is None:
                        request_context = {}
                    request_context["server_hostname"] = pm_self.hostname
                    return super()._new_pool(scheme, pm_self.pinned_ip, port, request_context)
                return super()._new_pool(scheme, host, port, request_context)

        self.poolmanager = PinnedPoolManager(self.hostname, self.pinned_ip, connections, maxsize, block=block, **pool_kwargs)


def _strip_repeated_lines(text: str, *, min_repeats: int = 3) -> str:
    lines = text.splitlines()
    normalized_counts: Dict[str, int] = {}
    for line in lines:
        normalized = re.sub(r"\s+", " ", line).strip().lower()
        if len(normalized) >= 4:
            normalized_counts[normalized] = normalized_counts.get(normalized, 0) + 1

    cleaned_lines = []
    for line in lines:
        normalized = re.sub(r"\s+", " ", line).strip().lower()
        if normalized_counts.get(normalized, 0) >= min_repeats:
            continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines)


def _clean_text(value: str) -> str:
    if not value:
        return ""
    text = value.replace("\x00", " ")
    for k, v in MOJIBAKE_REPLACEMENTS.items():
        text = text.replace(k, v)

    text = re.sub(r"(?i)^\s*halaman\s+\d+\s+(?:dari|of)\s+\d+\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"(?i)^\s*page\s+\d+\s+(?:of|dari)\s+\d+\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"(?i)^\s*(?:halaman|page)\s+\d+\s*$", "", text, flags=re.MULTILINE)
    text = _strip_repeated_lines(text)

    text = re.sub(r"(\d+)\s*\.\s*(\d{3})", r"\1.\2", text)
    text = re.sub(
        r"(\d+)\s*(?:\u00b0|degrees?)\s*C\b",
        lambda match: f"{match.group(1)}\u00b0C",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(
        r"(\d+(?:[.,]\d+)?)\s*(pa|mah|wh|w|v|ml|kg|g|cm|mm|m)\b",
        lambda match: f"{match.group(1)} {match.group(2)}",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(r"\b(\d+(?:[.,]\d+)?)\s+L\b", r"\1L", text)
    text = re.sub(r"(\d+)\s*°\s*C", r"\1°C", text)

    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _estimate_tokens(text: str) -> int:
    return max(1, len(text) // 4)


def _looks_like_table_line(stripped: str) -> bool:
    return ("|" in stripped and stripped.count("|") >= 2) or (
        "," in stripped and stripped.count(",") >= 3 and not stripped.endswith(".")
    )


def _looks_like_table_block(block_text: str) -> bool:
    lines = [line.strip() for line in block_text.splitlines() if line.strip()]
    return len(lines) >= 2 and all(_looks_like_table_line(line) for line in lines)


def _chunk_text(text: str, max_tokens: int = 1500, overlap_tokens: int = 150) -> List[str]:
    max_chars = max_tokens * 4
    overlap_chars = overlap_tokens * 4

    lines = text.splitlines()
    blocks: List[tuple[str | None, str]] = []
    current_block: List[str] = []
    current_header: str | None = None
    in_table = False

    header_re = re.compile(r"^\s*(#{1,4}\s+.+|[A-Z0-9\.\s]{3,60}:)\s*$")

    for line in lines:
        stripped = line.strip()
        if header_re.match(stripped) and not in_table:
            if current_block:
                blocks.append((current_header, "\n".join(current_block)))
                current_block = []
            current_header = stripped
            current_block.append(line)
            continue

        is_table_line = _looks_like_table_line(stripped)
        if is_table_line:
            if not in_table and current_block:
                blocks.append((current_header, "\n".join(current_block)))
                current_block = []
            in_table = True
            current_block.append(line)
        else:
            if in_table and current_block:
                delims = [l.count("|") if "|" in l else l.count(",") for l in current_block if l.strip()]
                if delims and len(set(delims)) > 1:
                    logger.warning(f"[TABLE_GUARDRAIL_WARNING] Inconsistent delimiter count across table rows: {set(delims)}")
                    current_block.insert(0, "[NEEDS_REVIEW: Struktur Baris/Kolom Tabel Tidak Konsisten]")
                blocks.append((current_header, "\n".join(current_block)))
                current_block = []
            in_table = False
            if not stripped:
                if current_block:
                    blocks.append((current_header, "\n".join(current_block)))
                    current_block = []
            else:
                current_block.append(line)

    if current_block:
        if in_table:
            delims = [l.count("|") if "|" in l else l.count(",") for l in current_block if l.strip()]
            if delims and len(set(delims)) > 1:
                logger.warning(f"[TABLE_GUARDRAIL_WARNING] Inconsistent delimiter count across table rows: {set(delims)}")
                current_block.insert(0, "[NEEDS_REVIEW: Struktur Baris/Kolom Tabel Tidak Konsisten]")
        blocks.append((current_header, "\n".join(current_block)))

    chunks: List[str] = []
    current_chunk = ""

    for hdr, block_text in blocks:
        b_str = block_text.strip()
        if not b_str:
            continue
        hdr_prefix = f"[{hdr}]\n" if (hdr and hdr not in b_str) else ""
        candidate = f"{current_chunk}\n\n{b_str}".strip() if current_chunk else f"{hdr_prefix}{b_str}".strip()

        if len(candidate) <= max_chars:
            current_chunk = candidate
        else:
            if current_chunk:
                chunks.append(current_chunk)
            if len(b_str) <= max_chars:
                current_chunk = f"{hdr_prefix}{b_str}".strip()
            elif _looks_like_table_block(b_str):
                chunks.append(f"{hdr_prefix}{b_str}".strip())
                current_chunk = ""
            else:
                sub_lines = b_str.splitlines()
                sub_chunk = ""
                for sub_line in sub_lines:
                    next_sub = f"{sub_chunk}\n{sub_line}".strip() if sub_chunk else f"{hdr_prefix}{sub_line}".strip()
                    if len(next_sub) <= max_chars:
                        sub_chunk = next_sub
                    else:
                        if sub_chunk:
                            chunks.append(sub_chunk)
                        sub_chunk = f"{hdr_prefix}{sub_line}".strip()
                current_chunk = sub_chunk

    if current_chunk:
        chunks.append(current_chunk)

    if len(chunks) <= 1:
        return chunks

    with_overlap: List[str] = [chunks[0]]
    for index in range(1, len(chunks)):
        if len(chunks[index]) > max_chars:
            with_overlap.append(chunks[index])
            continue
        prefix = chunks[index - 1][-overlap_chars:].strip()
        merged = f"{prefix}\n\n{chunks[index]}".strip()
        with_overlap.append(merged[: max_chars + overlap_chars])
    return with_overlap


class ReadableTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []
        self.skip_depth = 0

    def handle_starttag(self, tag: str, attrs: List[tuple[str, str | None]]) -> None:
        if tag in IGNORED_HTML_TAGS:
            self.skip_depth += 1
        if tag in {"p", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in IGNORED_HTML_TAGS and self.skip_depth > 0:
            self.skip_depth -= 1
        if tag in {"p", "li", "tr", "section", "article", "div"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self.skip_depth > 0:
            return
        text = data.strip()
        if text:
            self.parts.append(text)

    def get_text(self) -> str:
        return _clean_text(" ".join(self.parts))


def _resolve_and_validate_host(hostname: str) -> str:
    try:
        addresses = socket.getaddrinfo(hostname, None)
    except socket.gaierror as exc:
        raise HTTPException(status_code=400, detail=f"Hostname {hostname} tidak bisa di-resolve.") from exc

    resolved_ip = None
    for address in addresses:
        ip_str = address[4][0]
        try:
            ip = ipaddress.ip_address(ip_str)
        except ValueError:
            continue
        
        if (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_reserved
            or ip.is_unspecified
        ):
            raise HTTPException(
                status_code=400,
                detail=f"IP private/internal ({ip_str}) tidak diperbolehkan sebagai knowledge source."
            )
        
        if not resolved_ip:
            resolved_ip = ip_str

    if not resolved_ip:
        raise HTTPException(status_code=400, detail=f"Tidak ada IP publik yang valid ditemukan untuk {hostname}.")
    
    return resolved_ip


def _validate_public_web_url(url: str) -> str:
    parsed = urlparse(url.strip())
    if parsed.scheme not in {"http", "https"} or not parsed.netloc or not parsed.hostname:
        raise HTTPException(status_code=400, detail="URL web harus memakai http/https publik.")
    return parsed.geturl()


def _extract_web_page_text(url: str) -> str:
    current_url = _validate_public_web_url(url)
    headers = {
        "User-Agent": "SAISOKU-OMNIX-KnowledgeBot/1.0",
        "Accept": "text/html,application/xhtml+xml,text/plain;q=0.9,*/*;q=0.1",
    }

    for _ in range(4):
        parsed = urlparse(current_url)
        if not parsed.hostname:
            raise HTTPException(status_code=400, detail="URL redirect tidak valid.")
        
        ip = _resolve_and_validate_host(parsed.hostname)
        
        session = requests.Session()
        adapter = HostPinnedHTTPAdapter(parsed.hostname, ip)
        session.mount(f"http://{parsed.hostname}", adapter)
        session.mount(f"https://{parsed.hostname}", adapter)

        response = session.get(
            current_url,
            headers=headers,
            timeout=(5, 20),
            stream=True,
            allow_redirects=False,
        )
            
        if response.is_redirect or response.is_permanent_redirect:
            location = response.headers.get("location")
            if not location:
                raise HTTPException(status_code=400, detail="Redirect URL web tidak valid.")
            current_url = _validate_public_web_url(urljoin(current_url, location))
            continue

        if not response.ok:
            raise HTTPException(status_code=400, detail=f"Gagal membaca URL web: HTTP {response.status_code}.")

        content = bytearray()
        for chunk in response.iter_content(chunk_size=65536):
            if not chunk:
                continue
            content.extend(chunk)
            if len(content) > MAX_WEB_PAGE_BYTES:
                raise HTTPException(status_code=413, detail="Konten web terlalu besar untuk knowledge URL.")

        content_type = response.headers.get("content-type", "")
        encoding = response.encoding or "utf-8"
        raw_text = bytes(content).decode(encoding, errors="ignore")
        if "html" not in content_type.lower():
            return _clean_text(raw_text)

        parser = ReadableTextParser()
        parser.feed(raw_text)
        return parser.get_text()

    raise HTTPException(status_code=400, detail="URL web terlalu banyak redirect.")


def _extract_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PDF support requires pypdf to be installed",
        ) from exc

    reader = PdfReader(io.BytesIO(content))
    pages = [(page.extract_text() or "") for page in reader.pages]
    return "\n\n".join(pages)


def _extract_pdf_with_gemini_ocr(content: bytes) -> str:
    keys = _get_gemini_keys()
    if not keys:
        try:
            keys = [_gemini_api_key()]
        except Exception:
            keys = []
    candidate_models = chat_fallback_chain()
    encoded_pdf = base64.b64encode(content).decode("ascii")
    prompt = (
        "Transkripsikan teks dari PDF ini untuk knowledge base RAG. "
        "Baca juga halaman scan/gambar dengan OCR. "
        "Kembalikan hanya teks dokumen yang terbaca, pertahankan heading, tabel sederhana, "
        "nomor langkah, dan FAQ jika ada. Jangan membuat ringkasan atau menambah informasi."
    )
    last_err = ""
    for m in candidate_models:
        for key in keys:
            try:
                response = _HTTPX_CLIENT.post(
                    f"{GEMINI_API_BASE}/models/{m}:generateContent",
                    headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                    json={
                        "contents": [
                            {
                                "parts": [
                                    {"text": prompt},
                                    {
                                        "inlineData": {
                                            "mimeType": "application/pdf",
                                            "data": encoded_pdf,
                                        }
                                    },
                                ]
                            }
                        ],
                        "generationConfig": {"temperature": 0.1},
                    },
                )
                if response.status_code == 200:
                    candidates = response.json().get("candidates") or []
                    parts = (candidates[0].get("content", {}).get("parts") if candidates else []) or []
                    text = "\n".join(str(part.get("text", "")) for part in parts if part.get("text"))
                    cleaned = _clean_text(text)
                    if cleaned:
                        return cleaned
                else:
                    last_err = f"Model {m} HTTP {response.status_code}: {response.text[:200]}"
                    logger.warning(f"Gemini PDF OCR {last_err}")
            except Exception as exc:
                last_err = str(exc)
                logger.warning(f"Gemini PDF OCR request exception: {exc}")

    raise HTTPException(status_code=502, detail=f"Gagal mengekstrak teks dari PDF scan/OCR: {last_err}")


def _extract_docx(content: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="DOCX support requires python-docx to be installed",
        ) from exc

    document = Document(io.BytesIO(content))
    parts: List[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for index, table in enumerate(document.tables, start=1):
        rows = [
            [cell.text.strip() for cell in row.cells]
            for row in table.rows
            if any(cell.text.strip() for cell in row.cells)
        ]
        if rows:
            parts.append(f"### TABLE {index}\n{_markdown_table(rows)}")
    return "\n\n".join(parts)


def _markdown_table(rows: List[List[Any]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized_rows = [
        [str(cell).replace("\n", " ").strip() for cell in row] + [""] * (width - len(row))
        for row in rows
    ]
    header = normalized_rows[0]
    body = normalized_rows[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _dataframe_to_markdown(df: pd.DataFrame) -> str:
    cleaned_df = df.dropna(how="all").dropna(axis=1, how="all").fillna("")
    if cleaned_df.empty:
        return ""
    rows: List[List[Any]] = [list(cleaned_df.columns)]
    rows.extend(cleaned_df.astype(str).values.tolist())
    return _markdown_table(rows)


def _extract_spreadsheet(content: bytes, filename: str) -> str:
    file_obj = io.BytesIO(content)
    if filename.lower().endswith(".csv"):
        df = pd.read_csv(file_obj)
        if df.empty:
            return ""
        return _dataframe_to_markdown(df)
    else:
        excel_file = pd.ExcelFile(file_obj)
        sheet_texts: List[str] = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            if not df.empty:
                sheet_table = _dataframe_to_markdown(df)
                if sheet_table:
                    sheet_texts.append(f"### SHEET: {sheet_name}\n\n{sheet_table}")
        return "\n\n".join(sheet_texts)


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PPTX support requires python-pptx to be installed",
        ) from exc

    prs = Presentation(io.BytesIO(content))
    slide_texts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"### SLIDE {idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if hasattr(slide, "has_notes_slide") and slide.has_notes_slide and slide.notes_slide:
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame and notes_frame.text and notes_frame.text.strip():
                slide_lines.append(f"[Catatan Speaker Slide {idx}]: {notes_frame.text.strip()}")
        if len(slide_lines) > 1:
            slide_texts.append("\n".join(slide_lines))
    return "\n\n".join(slide_texts)


def _extract_image_with_gemini_ocr(content: bytes, mime_type: str = "image/png") -> str:
    keys = _get_gemini_keys()
    if not keys:
        try:
            keys = [_gemini_api_key()]
        except Exception:
            keys = []
    candidate_models = chat_fallback_chain()
    encoded_img = base64.b64encode(content).decode("ascii")
    prompt = (
        "Transkripsikan seluruh teks, tabel, spesifikasi produk, dan informasi penting "
        "yang ada di dalam gambar/foto ini untuk Knowledge Base RAG. "
        "Kembalikan hanya teks dokumen yang terbaca secara rapi tanpa membuat ringkasan opini atau menambah informasi."
    )
    last_err = ""
    for m in candidate_models:
        for key in keys:
            try:
                response = _HTTPX_CLIENT.post(
                    f"{GEMINI_API_BASE}/models/{m}:generateContent",
                    headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                    json={
                        "contents": [
                            {
                                "parts": [
                                    {"text": prompt},
                                    {
                                        "inlineData": {
                                            "mimeType": mime_type,
                                            "data": encoded_img,
                                        }
                                    },
                                ]
                            }
                        ],
                        "generationConfig": {"temperature": 0.1},
                    },
                )
                if response.status_code == 200:
                    candidates = response.json().get("candidates") or []
                    parts = (candidates[0].get("content", {}).get("parts") if candidates else []) or []
                    text = "\n".join(str(part.get("text", "")) for part in parts if part.get("text"))
                    cleaned = _clean_text(text)
                    if cleaned:
                        return cleaned
                else:
                    last_err = f"Model {m} HTTP {response.status_code}: {response.text[:200]}"
                    logger.warning(f"Gemini Image OCR {last_err}")
            except Exception as exc:
                last_err = str(exc)
                logger.warning(f"Gemini Image OCR exception: {exc}")

    raise HTTPException(status_code=502, detail=f"Gagal mengekstrak teks dari gambar: {last_err}")


def extract_document_text(content: bytes, filename: str, content_type: str | None) -> str:
    lower_name = filename.lower()
    if lower_name.endswith(".pdf") or content_type == "application/pdf":
        extracted_text = _clean_text(_extract_pdf(content))
        if len(extracted_text) >= MIN_EXTRACTED_TEXT_CHARS:
            return extracted_text
        return _extract_pdf_with_gemini_ocr(content)
    if lower_name.endswith(".docx"):
        return _clean_text(_extract_docx(content))
    if lower_name.endswith((".pptx", ".ppt")):
        return _clean_text(_extract_pptx(content))
    if lower_name.endswith((".xlsx", ".xls", ".csv")):
        return _clean_text(_extract_spreadsheet(content, filename))
    if lower_name.endswith((".jpg", ".jpeg", ".png", ".webp")) or (content_type and content_type.startswith("image/")):
        mtype = content_type if (content_type and "/" in content_type) else ("image/png" if lower_name.endswith(".png") else "image/jpeg")
        return _extract_image_with_gemini_ocr(content, mime_type=mtype)
    try:
        return _clean_text(content.decode("utf-8"))
    except UnicodeDecodeError:
        return _clean_text(content.decode("latin-1", errors="ignore"))


_CHUNK_CONTEXT_SYSTEM_INSTRUCTION = (
    "Anda membantu menyiapkan Knowledge Base RAG. Diberikan SATU DOKUMEN PENUH dan SATU "
    "POTONGAN (chunk) dari dokumen tersebut, tulis 1-2 kalimat singkat (maks 40 kata) dalam "
    "Bahasa Indonesia yang menjelaskan POSISI chunk ini di dalam dokumen — misalnya produk/model "
    "apa, bagian/section apa (spesifikasi, garansi, troubleshooting, indikator lampu, dst). "
    "JANGAN meringkas isi chunk. JANGAN menambah informasi yang tidak ada di dokumen. "
    "Jawab HANYA dengan kalimat context tersebut, tanpa embel-embel lain."
)


def _generate_chunk_context(document_text: str, chunk: str, document_title: str) -> str:
    keys = _get_gemini_keys()
    if not keys:
        return ""

    truncated_doc = document_text[:12000]
    candidate_models = chat_fallback_chain()
    prompt = (
        f"<dokumen judul=\"{document_title}\">\n{truncated_doc}\n</dokumen>\n\n"
        f"<chunk>\n{chunk[:2000]}\n</chunk>\n\n"
        "Tulis kalimat context untuk chunk di atas sesuai instruksi."
    )

    for m in candidate_models:
        for key in keys:
            try:
                response = _HTTPX_CLIENT.post(
                    f"{GEMINI_API_BASE}/models/{m}:generateContent",
                    headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                    json={
                        "systemInstruction": {"parts": [{"text": _CHUNK_CONTEXT_SYSTEM_INSTRUCTION}]},
                        "contents": [{"parts": [{"text": prompt}]}],
                        "generationConfig": {"temperature": 0.1, "maxOutputTokens": 100},
                    },
                )
                if response.status_code == 200:
                    candidates = response.json().get("candidates") or []
                    if candidates:
                        parts = candidates[0].get("content", {}).get("parts") or []
                        text = " ".join(str(p.get("text", "")) for p in parts if p.get("text"))
                        cleaned = _clean_text(text)
                        if cleaned:
                            return cleaned[:300]
                elif response.status_code == 429:
                    continue
            except Exception as exc:
                logger.warning(f"Chunk context generation failed: {exc}")

    return ""
