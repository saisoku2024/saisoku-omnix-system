import base64
from contextlib import contextmanager
import contextvars
from html.parser import HTMLParser
import ipaddress
import io
import logging
import os
import re
import socket
from typing import Any, Dict, List
from urllib.parse import urljoin, urlparse

import httpx
import numpy as np
import pandas as pd
import requests
from urllib3.util import connection
from fastapi import HTTPException, UploadFile, status

from app.core.supabase import supabase
from app.core.gemini_config import (
    resolve_chat_model,
    resolve_embedding_model,
    chat_fallback_chain,
    check_embedding_reindex_needed,
    DEFAULT_EMBEDDING_MODEL,
    EMBEDDING_DIMENSION,
)
from app.services.audit_log_service import AuditLogService
from app.services.storage_upload_service import (
    MAX_STORAGE_UPLOAD_SIZE_BYTES,
    validate_storage_upload,
)

logger = logging.getLogger(__name__)

GEMINI_API_BASE = "https://generativelanguage.googleapis.com/v1beta"
_HTTPX_CLIENT = httpx.Client(
    timeout=httpx.Timeout(60.0, connect=10.0),
    limits=httpx.Limits(max_keepalive_connections=20, max_connections=100),
)
MAX_KB_FILE_SIZE_BYTES = MAX_STORAGE_UPLOAD_SIZE_BYTES
MAX_WEB_PAGE_BYTES = 1 * 1024 * 1024
MIN_EXTRACTED_TEXT_CHARS = 20
IGNORED_HTML_TAGS = {"script", "style", "noscript", "svg", "nav", "header", "footer", "aside"}


def _check_and_flag_embedding_reindex_if_needed() -> None:
    if check_embedding_reindex_needed():
        logger.warning(
            f"[EMBEDDING_MISMATCH_WARNING] Environment GEMINI_EMBEDDING_MODEL differs from DEFAULT_EMBEDDING_MODEL ({DEFAULT_EMBEDDING_MODEL}). Flagging knowledge_documents as needs_reindex=True."
        )
        try:
            supabase.table("knowledge_documents").update({"needs_reindex": True}).eq("needs_reindex", False).execute()
        except Exception as exc:
            logger.warning(f"Failed to update knowledge_documents reindex flag: {exc}")

import urllib3
from requests.adapters import HTTPAdapter

class HostPinnedHTTPAdapter(HTTPAdapter):
    def __init__(self, hostname: str, pinned_ip: str, **kwargs):
        self.hostname = hostname
        self.pinned_ip = pinned_ip
        super().__init__(**kwargs)

    def init_poolmanager(self, connections, maxsize, block=False, **pool_kwargs):
        class PinnedPoolManager(urllib3.PoolManager):
            def __init__(pm_self, hostname, pinned_ip, *args, **kwargs):
                pm_self.hostname = hostname
                pm_self.pinned_ip = pinned_ip
                super().__init__(*args, **kwargs)

            def _new_pool(pm_self, scheme, host, port, request_context=None):
                if host == pm_self.hostname:
                    if request_context is None:
                        request_context = {}
                    request_context["server_hostname"] = pm_self.hostname
                    return super()._new_pool(scheme, pm_self.pinned_ip, port, request_context)
                return super()._new_pool(scheme, host, port, request_context)

        self.poolmanager = PinnedPoolManager(self.hostname, self.pinned_ip, connections, maxsize, block=block, **pool_kwargs)


_key_index = 0

def _get_gemini_keys() -> List[str]:
    keys_str = os.getenv("GEMINI_API_KEYS", "") or os.getenv("GEMINI_API_KEY", "")
    keys = [k.strip() for k in keys_str.split(",") if k.strip()]
    return keys

def _gemini_api_key() -> str:
    keys = _get_gemini_keys()
    if not keys:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="GEMINI_API_KEY is not configured",
        )
    global _key_index
    key = keys[_key_index % len(keys)]
    _key_index += 1
    return key



def _embedding_model() -> str:
    return resolve_embedding_model()


def _chat_model() -> str:
    return resolve_chat_model()


def _clean_text(value: str) -> str:
    text = value.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _estimate_tokens(text: str) -> int:
    return max(1, len(text) // 4)


def _chunk_text(text: str, max_tokens: int = 1500, overlap_tokens: int = 150) -> List[str]:
    """
    Chunking Table-Aware & Token-Aware.
    Mencegah pemotongan acak di tengah baris tabel/CSV dan menggunakan estimasi token (~4 karakter/token).
    """
    max_chars = max_tokens * 4
    overlap_chars = overlap_tokens * 4

    lines = text.splitlines()
    blocks: List[str] = []
    current_block: List[str] = []
    in_table = False

    for line in lines:
        stripped = line.strip()
        is_table_line = ("|" in stripped and stripped.count("|") >= 2) or ("," in stripped and stripped.count(",") >= 3 and not stripped.endswith("."))
        if is_table_line:
            if not in_table and current_block:
                blocks.append("\n".join(current_block))
                current_block = []
            in_table = True
            current_block.append(line)
        else:
            if in_table and current_block:
                # Check table delimiter consistency
                delims = [l.count("|") if "|" in l else l.count(",") for l in current_block if l.strip()]
                if delims and len(set(delims)) > 1:
                    logger.warning(f"[TABLE_GUARDRAIL_WARNING] Inconsistent delimiter count across table rows: {set(delims)}")
                    current_block.insert(0, "[NEEDS_REVIEW: Struktur Baris/Kolom Tabel Tidak Konsisten]")
                blocks.append("\n".join(current_block))
                current_block = []
            in_table = False
            if not stripped:
                if current_block:
                    blocks.append("\n".join(current_block))
                    current_block = []
            else:
                current_block.append(line)

    if current_block:
        if in_table:
            delims = [l.count("|") if "|" in l else l.count(",") for l in current_block if l.strip()]
            if delims and len(set(delims)) > 1:
                logger.warning(f"[TABLE_GUARDRAIL_WARNING] Inconsistent delimiter count across table rows: {set(delims)}")
                current_block.insert(0, "[NEEDS_REVIEW: Struktur Baris/Kolom Tabel Tidak Konsisten]")
        blocks.append("\n".join(current_block))

    chunks: List[str] = []
    current_chunk = ""

    for block in blocks:
        block_str = block.strip()
        if not block_str:
            continue
        next_chunk = f"{current_chunk}\n\n{block_str}".strip() if current_chunk else block_str
        if len(next_chunk) <= max_chars:
            current_chunk = next_chunk
        else:
            if current_chunk:
                chunks.append(current_chunk)
            if len(block_str) <= max_chars:
                current_chunk = block_str
            else:
                sub_lines = block_str.splitlines()
                sub_chunk = ""
                for sub_line in sub_lines:
                    next_sub = f"{sub_chunk}\n{sub_line}".strip() if sub_chunk else sub_line
                    if len(next_sub) <= max_chars:
                        sub_chunk = next_sub
                    else:
                        if sub_chunk:
                            chunks.append(sub_chunk)
                        sub_chunk = sub_line
                current_chunk = sub_chunk

    if current_chunk:
        chunks.append(current_chunk)

    if len(chunks) <= 1:
        return chunks

    with_overlap: List[str] = [chunks[0]]
    for index in range(1, len(chunks)):
        prefix = chunks[index - 1][-overlap_chars:].strip()
        merged = f"{prefix}\n\n{chunks[index]}".strip()
        with_overlap.append(merged[: max_chars + overlap_chars])
    return with_overlap


class ReadableTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []
        self.skip_depth = 0

    def handle_starttag(self, tag: str, attrs: List[tuple[str, str | None]]) -> None:
        if tag in IGNORED_HTML_TAGS:
            self.skip_depth += 1
        if tag in {"p", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in IGNORED_HTML_TAGS and self.skip_depth > 0:
            self.skip_depth -= 1
        if tag in {"p", "li", "tr", "section", "article", "div"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self.skip_depth > 0:
            return
        text = data.strip()
        if text:
            self.parts.append(text)

    def get_text(self) -> str:
        return _clean_text(" ".join(self.parts))


def _resolve_and_validate_host(hostname: str) -> str:
    try:
        addresses = socket.getaddrinfo(hostname, None)
    except socket.gaierror as exc:
        raise HTTPException(status_code=400, detail=f"Hostname {hostname} tidak bisa di-resolve.") from exc

    resolved_ip = None
    for address in addresses:
        ip_str = address[4][0]
        try:
            ip = ipaddress.ip_address(ip_str)
        except ValueError:
            continue
        
        if (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_reserved
            or ip.is_unspecified
        ):
            raise HTTPException(
                status_code=400,
                detail=f"IP private/internal ({ip_str}) tidak diperbolehkan sebagai knowledge source."
            )
        
        if not resolved_ip:
            resolved_ip = ip_str

    if not resolved_ip:
        raise HTTPException(status_code=400, detail=f"Tidak ada IP publik yang valid ditemukan untuk {hostname}.")
    
    return resolved_ip


def _validate_public_web_url(url: str) -> str:
    parsed = urlparse(url.strip())
    if parsed.scheme not in {"http", "https"} or not parsed.netloc or not parsed.hostname:
        raise HTTPException(status_code=400, detail="URL web harus memakai http/https publik.")
    return parsed.geturl()


def _extract_web_page_text(url: str) -> str:
    current_url = _validate_public_web_url(url)
    headers = {
        "User-Agent": "SAISOKU-OMNIX-KnowledgeBot/1.0",
        "Accept": "text/html,application/xhtml+xml,text/plain;q=0.9,*/*;q=0.1",
    }

    for _ in range(4):
        parsed = urlparse(current_url)
        if not parsed.hostname:
            raise HTTPException(status_code=400, detail="URL redirect tidak valid.")
        
        # Resolve and validate host to prevent SSRF
        ip = _resolve_and_validate_host(parsed.hostname)
        
        session = requests.Session()
        adapter = HostPinnedHTTPAdapter(parsed.hostname, ip)
        session.mount(f"http://{parsed.hostname}", adapter)
        session.mount(f"https://{parsed.hostname}", adapter)

        response = session.get(
            current_url,
            headers=headers,
            timeout=(5, 20),
            stream=True,
            allow_redirects=False,
        )
            
        if response.is_redirect or response.is_permanent_redirect:
            location = response.headers.get("location")
            if not location:
                raise HTTPException(status_code=400, detail="Redirect URL web tidak valid.")
            current_url = _validate_public_web_url(urljoin(current_url, location))
            continue

        if not response.ok:
            raise HTTPException(status_code=400, detail=f"Gagal membaca URL web: HTTP {response.status_code}.")

        content = bytearray()
        for chunk in response.iter_content(chunk_size=65536):
            if not chunk:
                continue
            content.extend(chunk)
            if len(content) > MAX_WEB_PAGE_BYTES:
                raise HTTPException(status_code=413, detail="Konten web terlalu besar untuk knowledge URL.")

        content_type = response.headers.get("content-type", "")
        encoding = response.encoding or "utf-8"
        raw_text = bytes(content).decode(encoding, errors="ignore")
        if "html" not in content_type.lower():
            return _clean_text(raw_text)

        parser = ReadableTextParser()
        parser.feed(raw_text)
        return parser.get_text()

    raise HTTPException(status_code=400, detail="URL web terlalu banyak redirect.")


def _extract_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PDF support requires pypdf to be installed",
        ) from exc

    reader = PdfReader(io.BytesIO(content))
    pages = [(page.extract_text() or "") for page in reader.pages]
    return "\n\n".join(pages)


def _extract_pdf_with_gemini_ocr(content: bytes) -> str:
    keys = _get_gemini_keys()
    if not keys:
        try:
            keys = [_gemini_api_key()]
        except Exception:
            keys = []
    candidate_models = chat_fallback_chain()
    encoded_pdf = base64.b64encode(content).decode("ascii")
    prompt = (
        "Transkripsikan teks dari PDF ini untuk knowledge base RAG. "
        "Baca juga halaman scan/gambar dengan OCR. "
        "Kembalikan hanya teks dokumen yang terbaca, pertahankan heading, tabel sederhana, "
        "nomor langkah, dan FAQ jika ada. Jangan membuat ringkasan atau menambah informasi."
    )
    last_err = ""
    for m in candidate_models:
        for key in keys:
            try:
                response = _HTTPX_CLIENT.post(
                    f"{GEMINI_API_BASE}/models/{m}:generateContent",
                    headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                    json={
                        "contents": [
                            {
                                "parts": [
                                    {"text": prompt},
                                    {
                                        "inlineData": {
                                            "mimeType": "application/pdf",
                                            "data": encoded_pdf,
                                        }
                                    },
                                ]
                            }
                        ],
                        "generationConfig": {"temperature": 0.1},
                    },
                )
                if response.status_code == 200:
                    candidates = response.json().get("candidates") or []
                    parts = (candidates[0].get("content", {}).get("parts") if candidates else []) or []
                    text = "\n".join(str(part.get("text", "")) for part in parts if part.get("text"))
                    cleaned = _clean_text(text)
                    if cleaned:
                        return cleaned
                else:
                    last_err = f"Model {m} HTTP {response.status_code}: {response.text[:200]}"
                    logger.warning(f"Gemini PDF OCR {last_err}")
            except Exception as exc:
                last_err = str(exc)
                logger.warning(f"Gemini PDF OCR request exception: {exc}")

    raise HTTPException(status_code=502, detail=f"Gagal mengekstrak teks dari PDF scan/OCR: {last_err}")


def _extract_docx(content: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="DOCX support requires python-docx to be installed",
        ) from exc

    document = Document(io.BytesIO(content))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_spreadsheet(content: bytes, filename: str) -> str:
    file_obj = io.BytesIO(content)
    if filename.lower().endswith(".csv"):
        df = pd.read_csv(file_obj)
        if df.empty:
            return ""
        return df.fillna("").astype(str).to_csv(index=False)
    else:
        excel_file = pd.ExcelFile(file_obj)
        sheet_texts: List[str] = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            if not df.empty:
                sheet_csv = df.fillna("").astype(str).to_csv(index=False)
                sheet_texts.append(f"### SHEET: {sheet_name}\n\n{sheet_csv}")
        return "\n\n".join(sheet_texts)



def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PPTX support requires python-pptx to be installed",
        ) from exc

    prs = Presentation(io.BytesIO(content))
    slide_texts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"### SLIDE {idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if hasattr(slide, "has_notes_slide") and slide.has_notes_slide and slide.notes_slide:
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame and notes_frame.text and notes_frame.text.strip():
                slide_lines.append(f"[Catatan Speaker Slide {idx}]: {notes_frame.text.strip()}")
        if len(slide_lines) > 1:
            slide_texts.append("\n".join(slide_lines))
    return "\n\n".join(slide_texts)


def _extract_image_with_gemini_ocr(content: bytes, mime_type: str = "image/png") -> str:
    keys = _get_gemini_keys()
    if not keys:
        try:
            keys = [_gemini_api_key()]
        except Exception:
            keys = []
    candidate_models = chat_fallback_chain()
    encoded_img = base64.b64encode(content).decode("ascii")
    prompt = (
        "Transkripsikan seluruh teks, tabel, spesifikasi produk, dan informasi penting "
        "yang ada di dalam gambar/foto ini untuk Knowledge Base RAG. "
        "Kembalikan hanya teks dokumen yang terbaca secara rapi tanpa membuat ringkasan opini atau menambah informasi."
    )
    last_err = ""
    for m in candidate_models:
        for key in keys:
            try:
                response = _HTTPX_CLIENT.post(
                    f"{GEMINI_API_BASE}/models/{m}:generateContent",
                    headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                    json={
                        "contents": [
                            {
                                "parts": [
                                    {"text": prompt},
                                    {
                                        "inlineData": {
                                            "mimeType": mime_type,
                                            "data": encoded_img,
                                        }
                                    },
                                ]
                            }
                        ],
                        "generationConfig": {"temperature": 0.1},
                    },
                )
                if response.status_code == 200:
                    candidates = response.json().get("candidates") or []
                    parts = (candidates[0].get("content", {}).get("parts") if candidates else []) or []
                    text = "\n".join(str(part.get("text", "")) for part in parts if part.get("text"))
                    cleaned = _clean_text(text)
                    if cleaned:
                        return cleaned
                else:
                    last_err = f"Model {m} HTTP {response.status_code}: {response.text[:200]}"
                    logger.warning(f"Gemini Image OCR {last_err}")
            except Exception as exc:
                last_err = str(exc)
                logger.warning(f"Gemini Image OCR exception: {exc}")

    raise HTTPException(status_code=502, detail=f"Gagal mengekstrak teks dari gambar: {last_err}")


def extract_document_text(content: bytes, filename: str, content_type: str | None) -> str:
    lower_name = filename.lower()
    if lower_name.endswith(".pdf") or content_type == "application/pdf":
        extracted_text = _clean_text(_extract_pdf(content))
        if len(extracted_text) >= MIN_EXTRACTED_TEXT_CHARS:
            return extracted_text
        return _extract_pdf_with_gemini_ocr(content)
    if lower_name.endswith(".docx"):
        return _clean_text(_extract_docx(content))
    if lower_name.endswith((".pptx", ".ppt")):
        return _clean_text(_extract_pptx(content))
    if lower_name.endswith((".xlsx", ".xls", ".csv")):
        return _clean_text(_extract_spreadsheet(content, filename))
    if lower_name.endswith((".jpg", ".jpeg", ".png", ".webp")) or (content_type and content_type.startswith("image/")):
        mtype = content_type if (content_type and "/" in content_type) else ("image/png" if lower_name.endswith(".png") else "image/jpeg")
        return _extract_image_with_gemini_ocr(content, mime_type=mtype)
    try:
        return _clean_text(content.decode("utf-8"))
    except UnicodeDecodeError:
        return _clean_text(content.decode("latin-1", errors="ignore"))


def _normalize_l2(values: List[float], expected_dim: int = EMBEDDING_DIMENSION) -> List[float]:
    if not values or len(values) != expected_dim:
        return [0.0] * expected_dim
    arr = np.array(values, dtype=np.float32)
    norm = np.linalg.norm(arr)
    if norm > 0:
        arr = arr / norm
    return arr.tolist()


def _embed_text(text: str, *, title: str | None = None, is_query: bool = False) -> List[float]:
    keys = _get_gemini_keys()
    if not keys:
        try:
            keys = [_gemini_api_key()]
        except Exception:
            keys = []
    model = _embedding_model()
    prefix = "task: search result | query: " if is_query else f"title: {title or 'none'} | text: "
    payload = {
        "content": {"parts": [{"text": f"{prefix}{text}"}]},
        "output_dimensionality": EMBEDDING_DIMENSION,
    }

    for key in keys:
        try:
            response = _HTTPX_CLIENT.post(
                f"{GEMINI_API_BASE}/models/{model}:embedContent",
                headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                json=payload,
            )
            if response.status_code == 200:
                values = response.json().get("embedding", {}).get("values")
                if isinstance(values, list) and len(values) == EMBEDDING_DIMENSION:
                    return _normalize_l2([float(v) for v in values])
            elif response.status_code == 429:
                logger.warning("Gemini embedding key hit rate limit (429). Trying next key...")
                continue
        except Exception as exc:
            logger.warning(f"Gemini embedding exception: {exc}")

    # Fallback to zero vector if embedding fails or rate limited
    logger.warning("All Gemini embedding keys failed. Returning 768-dim zero vector fallback.")
    return [0.0] * EMBEDDING_DIMENSION



def _embed_texts(texts: List[str], *, title: str | None = None) -> List[List[float]]:
    if not texts:
        return []
    keys = _get_gemini_keys()
    if not keys:
        keys = [_gemini_api_key()]
    model = _embedding_model()
    model_name = model if model.startswith("models/") else f"models/{model}"
    
    # Max batch size for Gemini batchEmbedContents is 100
    batch_size = 100
    all_embeddings = []
    
    for i in range(0, len(texts), batch_size):
        chunk_batch = texts[i : i + batch_size]
        requests_payload = []
        for text in chunk_batch:
            prefix = f"title: {title or 'none'} | text: "
            requests_payload.append({
                "model": model_name,
                "content": {"parts": [{"text": f"{prefix}{text}"}]},
                "output_dimensionality": EMBEDDING_DIMENSION,
            })
            
        payload = {"requests": requests_payload}
        response = None
        last_error = ""
        # Rotasi API key: satu key kena rate limit (429) jangan bikin seluruh ingest gagal.
        for key in keys:
            response = _HTTPX_CLIENT.post(
                f"{GEMINI_API_BASE}/models/{model}:batchEmbedContents",
                headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                json=payload,
            )
            if response.status_code == 200:
                break
            last_error = response.text[:300]
            if response.status_code == 429:
                logger.warning("Gemini batch embedding key rate limited (429). Trying next key...")
                continue
            break  # non-429 error, no point retrying other keys with same bad payload
        if not response or response.status_code != 200:
            raise HTTPException(
                status_code=502,
                detail=f"Gemini batch embedding request failed: {last_error}",
            )
        
        embeddings_data = response.json().get("embeddings") or []
        if len(embeddings_data) != len(chunk_batch):
            raise HTTPException(
                status_code=502,
                detail="Gemini batch embedding returned mismatching number of embeddings"
            )
            
        for emb in embeddings_data:
            values = emb.get("values")
            if not isinstance(values, list) or len(values) != EMBEDDING_DIMENSION:
                raise HTTPException(status_code=502, detail="Gemini batch embedding response is invalid")
            all_embeddings.append(_normalize_l2([float(v) for v in values]))
            
    return all_embeddings


def _vector_literal(values: List[float]) -> str:
    return "[" + ",".join(f"{v:.8f}" for v in values) + "]"


_KB_SYSTEM_INSTRUCTION = (
    "Anda adalah AI Knowledge Base untuk SAISOKU OMNIX. "
    "Jawab dalam Bahasa Indonesia yang ringkas, rapi, dan HANYA berdasarkan konteks yang diberikan. "
    "DILARANG MENAMPILKAN CATATAN INTERNAL, VERIFIKASI BARIS/KOLOM, ATAU PROSES BERPIKIR DRAFT DI DALAM JAWABAN. "
    "JIKA TERDETEKSI KETIDAK-KONSISTENAN / PERBEDAAN DATA ANTAR-DOKUMEN DI DALAM KONTEKS (misalnya Dokumen A menyebut angka/nilai X dan Dokumen B menyebut angka/nilai Y): "
    "WAJIB gunakan PRINSIP TRANSPARANSI dengan menyajikan: "
    "1. Menyebutkan nilai/informasi pada dokumen versi lama/terdahulu (misal: Dokumen A - 290 menit). "
    "2. Menyebutkan nilai/informasi pada dokumen versi baru/terbaru (misal: Dokumen B - 240 menit). "
    "3. Menyampaikan rekomendasi rujukan utama secara tegas berdasarkan publikasi dokumen yang paling baru. "
    "Jika pertanyaan berupa PERBANDINGAN / PERBEDAAAN (seperti 'beda X dan Y', 'vs', 'perbandingan'), WAJIB sajikan dengan: "
    "1. Tabel Markdown perbandingan fitur yang berbeda. "
    "2. Poin-poin persamaan fitur. "
    "3. Ringkasan kesimpulan singkat. "
    "Untuk pertanyaan spesifikasi produk atau informasi umum, sajikan poin spesifikasi secara terstruktur lalu AKHIRI DENGAN '📌 Ringkasan Keunggulan / Kesimpulan' singkat di bagian bawah. "
    "Jika konteks tidak cukup untuk menjawab, katakan dengan jelas bahwa knowledge base belum punya "
    "informasi yang cukup, jangan menebak atau mengarang."
)


def _generate_answer(question: str, sources: List[Dict[str, Any]]) -> str:
    context = "\n\n".join(
        f"[Source {idx + 1}: {source['title']}]\n{source['content']}"
        for idx, source in enumerate(sources)
    )
    prompt = f"KONTEKS:\n{context}\n\nPERTANYAAN:\n{question}"

    # Try Gemini API Keys
    keys = _get_gemini_keys()
    if not keys:
        try:
            keys = [_gemini_api_key()]
        except Exception:
            keys = []
    candidate_models = chat_fallback_chain()

    for m in candidate_models:
        for key in keys:
            try:
                response = _HTTPX_CLIENT.post(
                    f"{GEMINI_API_BASE}/models/{m}:generateContent",
                    headers={"x-goog-api-key": key, "Content-Type": "application/json"},
                    json={
                        "systemInstruction": {"parts": [{"text": _KB_SYSTEM_INSTRUCTION}]},
                        "contents": [{"parts": [{"text": prompt}]}],
                        "generationConfig": {"temperature": 0.2, "maxOutputTokens": 2048},
                    },
                )
                if response.status_code == 200:
                    candidates = response.json().get("candidates") or []
                    if candidates:
                        content = candidates[0].get("content", {})
                        parts = content.get("parts") or []
                        final_parts = [p for p in parts if not p.get("thought") and not p.get("thoughtSignature")]
                        if not final_parts:
                            final_parts = [p for p in parts if not p.get("thought")]
                        if not final_parts and parts:
                            final_parts = [parts[-1]]
                        text = "\n".join(str(part.get("text", "")) for part in final_parts if part.get("text"))
                        if text.strip():
                            return text.strip()
                elif response.status_code == 429:
                    logger.warning(f"Gemini API key ({key[:6]}...) rate limited on {m}. Trying next key...")
                    continue
                else:
                    logger.warning(f"Gemini model {m} call failed HTTP {response.status_code}: {response.text[:150]}")
            except Exception as exc:
                logger.warning(f"Gemini LLM request failed for model {m}: {exc}")

    # Fallback if LLM API key fails, returns 429, or unavailable
    return "Berikut informasi yang ditemukan di Knowledge Base SOP:\n\n" + context





class KnowledgeService:
    @staticmethod
    def list_documents(limit: int = 50, offset: int = 0) -> Dict[str, Any]:
        safe_limit = max(1, min(limit, 200))
        safe_offset = max(0, offset)
        end_range = safe_offset + safe_limit - 1

        try:
            res = (
                supabase.table("knowledge_documents")
                .select("id,title,source_file,mime_type,status,chunk_count,created_by,error_summary,storage_bucket,storage_path,file_size,created_at,updated_at", count="exact")
                .order("created_at", desc=True)
                .range(safe_offset, end_range)
                .execute()
            )
        except Exception as exc:
            if not any(key in str(exc).lower() for key in ["storage_bucket", "storage_path", "file_size"]):
                raise
            res = (
                supabase.table("knowledge_documents")
                .select("id,title,source_file,mime_type,status,chunk_count,created_by,error_summary,created_at,updated_at", count="exact")
                .order("created_at", desc=True)
                .range(safe_offset, end_range)
                .execute()
            )
        documents = res.data or []
        total_count = res.count if getattr(res, "count", None) is not None else len(documents)
        return {
            "total": total_count,
            "limit": safe_limit,
            "offset": safe_offset,
            "documents": documents,
        }

    @staticmethod
    async def prepare_upload(file: UploadFile, title: str | None, user_email: str = "admin@omnix.com") -> Dict[str, Any]:
        content = await file.read()
        validate_storage_upload("knowledge", file.filename or "", len(content))

        document_title = (title or file.filename or "Untitled Knowledge Document").strip()
        doc_res = (
            supabase.table("knowledge_documents")
            .insert(
                {
                    "title": document_title,
                    "source_file": file.filename,
                    "mime_type": file.content_type,
                    "status": "processing",
                    "created_by": user_email,
                }
            )
            .execute()
        )
        document = (doc_res.data or [None])[0]
        if not document:
            raise HTTPException(status_code=500, detail="Gagal membuat knowledge document.")

        return {
            "success": True,
            "document_id": document["id"],
            "title": document_title,
            "status": "processing",
            "source_file": file.filename,
            "content_type": file.content_type,
            "content": content,
        }

    @staticmethod
    def prepare_storage_upload(
        filename: str,
        title: str | None,
        content_type: str | None,
        storage_bucket: str,
        storage_path: str,
        file_size: int,
        user_email: str = "admin@omnix.com",
    ) -> Dict[str, Any]:
        document_title = (title or filename or "Untitled Knowledge Document").strip()
        payload = {
            "title": document_title,
            "source_file": filename,
            "mime_type": content_type,
            "status": "processing",
            "created_by": user_email,
            "storage_bucket": storage_bucket,
            "storage_path": storage_path,
            "file_size": file_size,
        }
        try:
            doc_res = supabase.table("knowledge_documents").insert(payload).execute()
        except Exception as exc:
            if not any(key in str(exc).lower() for key in ["storage_bucket", "storage_path", "file_size"]):
                raise
            fallback_payload = {
                "title": document_title,
                "source_file": f"{storage_bucket}/{storage_path}",
                "mime_type": content_type,
                "status": "processing",
                "created_by": user_email,
            }
            doc_res = supabase.table("knowledge_documents").insert(fallback_payload).execute()
        document = (doc_res.data or [None])[0]
        if not document:
            raise HTTPException(status_code=500, detail="Gagal membuat knowledge document.")

        return {
            "success": True,
            "document_id": document["id"],
            "title": document_title,
            "status": "processing",
            "source_file": filename,
            "content_type": content_type,
            "storage_bucket": storage_bucket,
            "storage_path": storage_path,
        }

    @staticmethod
    def _process_text_content(
        document_id: str,
        document_title: str,
        text: str,
        user_email: str = "admin@omnix.com",
    ) -> None:
        try:
            text = _clean_text(text)
            if len(text) < MIN_EXTRACTED_TEXT_CHARS:
                raise HTTPException(status_code=400, detail="Dokumen terlalu kosong untuk diproses sebagai knowledge base.")

            chunks = _chunk_text(text)
            if not chunks:
                raise HTTPException(status_code=400, detail="Dokumen tidak menghasilkan chunk knowledge base.")

            embeddings = _embed_texts(chunks, title=document_title)
            rows = []
            for index, chunk in enumerate(chunks):
                rows.append(
                    {
                        "document_id": document_id,
                        "chunk_index": index,
                        "title": document_title,
                        "content": chunk,
                        "token_estimate": _estimate_tokens(chunk),
                        "embedding": _vector_literal(embeddings[index]),
                    }
                )
            supabase.table("knowledge_chunks").insert(rows).execute()
            (
                supabase.table("knowledge_documents")
                .update({"status": "ready", "chunk_count": len(rows), "error_summary": None})
                .eq("id", document_id)
                .execute()
            )
            AuditLogService.log(
                action="KNOWLEDGE_UPLOAD",
                resource="knowledge_documents",
                user_email=user_email,
                user_role="super_admin",
                details={"document_id": document_id, "title": document_title, "chunks": len(rows)},
            )
        except Exception as exc:
            logger.error("Knowledge ingestion failed", exc_info=True)
            error_summary = exc.detail if isinstance(exc, HTTPException) else str(exc)
            (
                supabase.table("knowledge_documents")
                .update({"status": "failed", "error_summary": str(error_summary)[:500]})
                .eq("id", document_id)
                .execute()
            )

    @staticmethod
    def process_upload(
        document_id: str,
        content: bytes,
        filename: str | None,
        content_type: str | None,
        document_title: str,
        user_email: str = "admin@omnix.com",
    ) -> None:
        try:
            try:
                text = extract_document_text(content, filename or document_title, content_type)
            except HTTPException:
                raise
            except Exception as exc:
                raise HTTPException(
                    status_code=400,
                    detail=f"Gagal membaca dokumen knowledge base: {str(exc)[:300]}",
                ) from exc
            KnowledgeService._process_text_content(document_id, document_title, text, user_email)
        except Exception as exc:
            logger.error("Knowledge ingestion failed", exc_info=True)
            error_summary = exc.detail if isinstance(exc, HTTPException) else str(exc)
            (
                supabase.table("knowledge_documents")
                .update({"status": "failed", "error_summary": str(error_summary)[:500]})
                .eq("id", document_id)
                .execute()
            )

    @staticmethod
    def prepare_manual_text(title: str, text: str, user_email: str = "admin@omnix.com") -> Dict[str, Any]:
        document_title = title.strip()
        cleaned_text = _clean_text(text)
        if len(document_title) < 3:
            raise HTTPException(status_code=400, detail="Judul knowledge manual minimal 3 karakter.")
        if len(cleaned_text) < MIN_EXTRACTED_TEXT_CHARS:
            raise HTTPException(status_code=400, detail="Teks knowledge manual terlalu pendek.")

        doc_res = (
            supabase.table("knowledge_documents")
            .insert(
                {
                    "title": document_title,
                    "source_file": "manual:text",
                    "mime_type": "text/plain",
                    "status": "processing",
                    "created_by": user_email,
                }
            )
            .execute()
        )
        document = (doc_res.data or [None])[0]
        if not document:
            raise HTTPException(status_code=500, detail="Gagal membuat manual knowledge document.")

        return {
            "success": True,
            "document_id": document["id"],
            "title": document_title,
            "status": "processing",
            "text": cleaned_text,
        }

    @staticmethod
    def process_manual_text(
        document_id: str,
        document_title: str,
        text: str,
        user_email: str = "admin@omnix.com",
    ) -> None:
        KnowledgeService._process_text_content(document_id, document_title, text, user_email)

    @staticmethod
    def prepare_web_url(url: str, title: str | None = None, user_email: str = "admin@omnix.com") -> Dict[str, Any]:
        source_url = _validate_public_web_url(url)
        parsed = urlparse(source_url)
        document_title = (title or parsed.netloc).strip()
        if len(document_title) < 3:
            raise HTTPException(status_code=400, detail="Judul web knowledge minimal 3 karakter.")

        doc_res = (
            supabase.table("knowledge_documents")
            .insert(
                {
                    "title": document_title,
                    "source_file": source_url,
                    "mime_type": "text/html",
                    "status": "processing",
                    "created_by": user_email,
                }
            )
            .execute()
        )
        document = (doc_res.data or [None])[0]
        if not document:
            raise HTTPException(status_code=500, detail="Gagal membuat web knowledge document.")

        return {
            "success": True,
            "document_id": document["id"],
            "title": document_title,
            "status": "processing",
            "url": source_url,
        }

    @staticmethod
    def process_web_url(
        document_id: str,
        document_title: str,
        url: str,
        user_email: str = "admin@omnix.com",
    ) -> None:
        try:
            text = _extract_web_page_text(url)
            KnowledgeService._process_text_content(document_id, document_title, text, user_email)
        except Exception as exc:
            logger.error("Knowledge web ingestion failed", exc_info=True)
            error_summary = exc.detail if isinstance(exc, HTTPException) else str(exc)
            (
                supabase.table("knowledge_documents")
                .update({"status": "failed", "error_summary": str(error_summary)[:500]})
                .eq("id", document_id)
                .execute()
            )

    @staticmethod
    def query(question: str, match_count: int = 6) -> Dict[str, Any]:
        _check_and_flag_embedding_reindex_if_needed()
        cleaned_question = question.strip()
        if len(cleaned_question) < 3:
            raise HTTPException(status_code=400, detail="Pertanyaan terlalu pendek.")

        sources: List[Dict[str, Any]] = []

        # 1. Vector Search (if Gemini Embedding API key is valid & non-zero)
        try:
            embedding = _embed_text(cleaned_question, is_query=True)
            if any(v != 0.0 for v in embedding):
                res = (
                    supabase.rpc(
                        "match_knowledge_chunks",
                        {
                            "query_embedding": _vector_literal(embedding),
                            "match_count": match_count,
                        },
                    )
                    .execute()
                )
                raw_sources = res.data or []
                sources = [s for s in raw_sources if float(s.get("similarity") or 0) >= 0.5]
        except Exception as exc:
            logger.warning(f"Vector search embedding failed or unconfigured: {exc}")

        stop_words = {
            "apa", "yang", "dan", "atau", "dengan", "pada", "untuk", "dari", "ke", "ini", "itu",
            "adalah", "bisa", "bagaimana", "mengapa", "apakah", "berapa", "saya", "tanya", "sesuai",
            "sisi", "dokumen", "perbedaan", "spek", "spesifikasi", "info", "informasi", "fitur",
            "detail", "tolong", "minta", "kasih", "tahu", "jelaskan",
            # Bentuk kolokial yang tadinya lolos filter dan bikin pencarian jadi generik
            "beda", "bedanya", "vs", "versus", "sama", "kayak", "gimana", "kenapa", "gak", "nggak",
            "antara", "dibanding", "dibandingkan", "lebih",
        }
        words = [re.sub(r"[^\w\.]", "", w).strip() for w in cleaned_question.split()]

        def _is_meaningful_keyword(w: str) -> bool:
            if not w or w.lower() in stop_words:
                return False
            if len(w) >= 3:
                return True
            return len(w) == 2 and any(c.isdigit() for c in w)

        raw_keywords = [w for w in words if _is_meaningful_keyword(w)]
        seen_kw: set = set()
        keywords = []
        for w in raw_keywords:
            lw = w.lower()
            if lw not in seen_kw:
                seen_kw.add(lw)
                keywords.append(w)
        keywords.sort(key=lambda w: (not any(c.isdigit() for c in w), len(w)))

        # Entitas produk spesifik (misal: "Y1", "S9", "T80")
        product_code_keywords = [w for w in keywords if any(c.isdigit() for c in w)]

        # Jika vector search menghasilkan sumber tetapi TIDAK SATUPUN sumber memuat kode produk spesifik
        # yang ditanyakan user, maka vector search terjebak di dokumen generik (misal: Troubleshooting).
        # Kita saring sumber yang tidak cocok agar keyword search mengambil dokumen produk yang tepat.
        if sources and product_code_keywords:
            matching_sources = [
                s for s in sources
                if any(pk.lower() in (s.get("content") or "").lower() or pk.lower() in (s.get("title") or "").lower() for pk in product_code_keywords)
            ]
            if matching_sources:
                sources = matching_sources
            else:
                logger.info(f"Vector search returned generic chunks without product codes {product_code_keywords}. Falling back to keyword search.")
                sources = []

        # 2. Keyword Search Fallback if vector search yields insufficient relevant sources
        if len(sources) < match_count:

            if keywords:
                existing_ids = {s.get("chunk_id") for s in sources if s.get("chunk_id")}

                # Tier 1: AND semua top-3 keyword — presisi tinggi buat pertanyaan satu topik
                query_builder = supabase.table("knowledge_chunks").select("id, document_id, title, content, chunk_index")
                for kw in keywords[:3]:
                    query_builder = query_builder.ilike("content", f"%{kw}%")
                kw_res = query_builder.limit(match_count).execute()
                kw_chunks = kw_res.data or []

                # Tier 2: OR top keyword-keyword spesifik — buat pertanyaan perbandingan
                # ("Y1 vs Y1 Pro") di mana tiap istilah ada di chunk/dokumen berbeda,
                # bukan nyampur di satu chunk yang sama.
                if not kw_chunks and len(keywords) > 1:
                    or_filter = ",".join(f"content.ilike.%{kw}%" for kw in keywords[:4])
                    kw_res = (
                        supabase.table("knowledge_chunks")
                        .select("id, document_id, title, content, chunk_index")
                        .or_(or_filter)
                        .limit(match_count * 2)
                        .execute()
                    )
                    kw_chunks = kw_res.data or []

                # Tier 3: fallback ke keyword TERSPESIFIK (bukan kata pertama di kalimat)
                if not kw_chunks:
                    query_builder = supabase.table("knowledge_chunks").select("id, document_id, title, content, chunk_index")
                    kw_res = query_builder.ilike("content", f"%{keywords[0]}%").limit(match_count).execute()
                    kw_chunks = kw_res.data or []

                for kc in kw_chunks:
                    chunk_id = kc.get("id")
                    if chunk_id and chunk_id not in existing_ids:
                        sources.append({
                            "chunk_id": chunk_id,
                            "document_id": kc.get("document_id"),
                            "title": kc.get("title"),
                            "content": kc.get("content"),
                            "chunk_index": kc.get("chunk_index"),
                            "similarity": 0.95,
                        })
                        existing_ids.add(chunk_id)



        if not sources:
            return {
                "answer": "Knowledge base belum punya informasi yang cukup untuk menjawab pertanyaan ini.",
                "sources": [],
            }

        answer = _generate_answer(cleaned_question, sources[:match_count])
        AuditLogService.log(
            action="KNOWLEDGE_QUERY",
            resource="knowledge_chunks",
            details={"question": cleaned_question, "source_count": len(sources[:match_count])},
        )
        return {
            "answer": answer,
            "sources": [
                {
                    "chunk_id": source.get("chunk_id"),
                    "document_id": source.get("document_id"),
                    "title": source.get("title"),
                    "content": source.get("content"),
                    "chunk_index": source.get("chunk_index"),
                    "similarity": source.get("similarity"),
                }
                for source in sources[:match_count]
            ],
        }

